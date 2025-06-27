from pptx import Presentation
from pptx.util import Inches, Pt
import math
from pptx.enum.text import PP_ALIGN
import os
import json

# 设置字体
def set_font(placeholder, font_size, font_name="Microsoft YaHei"):
    """
    设置段落的字体和字号
    """
    for paragraph in placeholder.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.name = font_name
            run.font.bold = None  # 可根据需要调整

# 适应文本框
def fit_text_to_box(shape, text, max_font_size=32, min_font_size=14):
    # 先设置最大字号
    shape.text = text
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(max_font_size)
    # 简单估算：如果字符数过多，缩小字号
    char_count = len(text)
    if char_count > 100:
        font_size = max(min_font_size, max_font_size - (char_count - 100) // 5)
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)

# 添加封面
def add_cover_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    set_font(slide.shapes.title, 48)
    set_font(slide.placeholders[1], 28)
    # 可加logo、背景色等美化

# 添加目录
def add_toc_slide(prs, toc_list):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "目录"
    content_shape = slide.placeholders[1]
    content_shape.text = ""
    for idx, item in enumerate(toc_list, 1):
        p = content_shape.text_frame.add_paragraph()
        p.text = f"{idx}. {item}"
        p.level = 0
    set_font(slide.shapes.title, 36)
    set_font(content_shape, 24)

# 添加章节
def add_section_slide(prs, section_title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白页
    left = top = Inches(2)
    width = height = Inches(6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = section_title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# 添加小标题
def add_content_slide(prs, title, content, para_ids, mapping, page_num, subtitle_font_size=20):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shape = slide.shapes.title
    shape.text = title
    # 设置主标题字体
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(subtitle_font_size if "——" in title else 32)
            run.font.name = "Microsoft YaHei"
    content_shape = slide.placeholders[1]
    fit_text_to_box(content_shape, content)
    mapping[str(page_num)] = para_ids


# 分割内容
def split_content(text, max_chars=300):
    # 按最大字符数分割
    return [text[i:i+max_chars] for i in range(0, len(text), max_chars)]

# 添加多内容
def add_multi_content_slide(prs, title, content):
    parts = split_content(content)
    for part in parts:
        add_content_slide(prs, title, part)

# 分组内容
def group_content_items(items, max_items_per_slide=5, max_chars_per_slide=300):
    slides = []
    current_slide = []
    current_chars = 0
    for item in items:
        if (len(current_slide) >= max_items_per_slide) or (current_chars + len(item['text']) > max_chars_per_slide):
            slides.append(current_slide)
            current_slide = []
            current_chars = 0
        current_slide.append(item)
        current_chars += len(item['text'])
    if current_slide:
        slides.append(current_slide)
    return slides

# 获取幻灯片标题
def get_slide_title(section, subsection):
    if subsection:
        return f"{section} —— {subsection}"
    else:
        return section
# 生成PPT
def generate_ppt(doc_id, structure):
    structure = get_structure_for_ppt(doc_id)
    prs = Presentation()
    mapping = {}
    page_num = 1

    # 1. 封面
    add_cover_slide(prs, "智绘PPT", "专业内容自动生成")
    mapping[str(page_num)] = ["cover"]  # 封面不对应原文
    page_num += 1

    # 2. 目录页（只列出一级标题）
    toc_titles = [item['text'] for item in structure if item.get('type') == 'title' and item.get('level') == 1]
    add_toc_slide(prs, toc_titles)
    mapping[str(page_num)] = []
    page_num += 1

    # 3. 内容
    current_section = None
    current_section_id = None
    current_subsection = None
    content_items = []

    def flush_content():
        nonlocal page_num, content_items, current_section, current_subsection
        if current_section and content_items:
            slides = group_content_items(content_items)
            for slide_items in slides:
                slide_title = get_slide_title(current_section, current_subsection)
                content = "\n".join([f"• {i['text']}" for i in slide_items if i.get('type') == 'paragraph' or i.get('type') == 'subtitle'])
                para_ids = [i['id'] for i in slide_items if i.get('type') == 'paragraph']
                add_content_slide(prs, slide_title, content, para_ids, mapping, page_num, subtitle_font_size=20)
                page_num += 1
            content_items = []

    for item in structure:
        if item.get('type') == 'title' and item.get('level') == 1:
            flush_content()
            current_section = item['text']
            current_section_id = item['id']
            current_subsection = None
            add_section_slide(prs, current_section)
            mapping[str(page_num)] = [current_section_id]
            page_num += 1
        elif item.get('type') == 'subtitle' and item.get('level') == 2:
            flush_content()
            current_subsection = item['text']
        else:
            content_items.append(item)
    # 处理最后一节
    flush_content()

    ppt_path = f"output/{doc_id}.pptx"
    prs.save(ppt_path)
    # 保存映射表
    with open(f"output/{doc_id}_mapping.json", "w", encoding="utf-8") as f:
        json.dump(mapping, f, ensure_ascii=False)
    return ppt_path
# 生成PPT
def generate_ppt_with_template(doc_id, structure, template='none'):
    if template == 'none':
        prs = Presentation()  # 空白PPT
        # 用默认代码生成各页
        # ...
    else:
        template_path = f"server/templates/{template}.pptx"
        config_path = f"server/template_configs/{template}.json"
        prs = Presentation(template_path)
        with open(config_path, encoding='utf-8') as f:
            config = json.load(f)
        mapping = {}
        page_num = 1

        # 标题页
        fill_title_slide(prs.slides[0], structure, config['title'])
        mapping[str(page_num)] = ['title']  # 假设title页只对应title
        page_num += 1

        # 目录页
        fill_toc_slide(prs.slides[1], structure, config['toc'])
        mapping[str(page_num)] = []  # 目录页不映射原文
        page_num += 1

        # 章节页、内容页
        # 遍历结构化内容，遇到章节/内容时分别填充，并记录映射
        # 例如：
        for section in sections:
            fill_section_slide(prs.slides[page_num], section['text'], config['section'])
            mapping[str(page_num)] = [section['id']]
            page_num += 1
            for para in section['paras']:
                fill_content_slide(prs.slides[page_num], para['title'], para['content'], config['content'])
                mapping[str(page_num)] = [para['id']]
                page_num += 1

        # 结束页
        fill_end_slide(prs.slides[page_num], "感谢聆听", config['end'])
        mapping[str(page_num)] = ['end']

        # 保存映射
        with open(f'output/{doc_id}_mapping.json', 'w', encoding='utf-8') as f:
            json.dump(mapping, f, ensure_ascii=False)
        ppt_path = f"output/{doc_id}.pptx"
        os.makedirs(os.path.dirname(ppt_path), exist_ok=True)
        # 保存PPT
        prs.save(ppt_path)
    return f'output/{doc_id}.pptx'

def get_structure_for_ppt(doc_id):
    middle_path = f"output/{doc_id}_middle.json"
    structure_path = f"output/{doc_id}_structure.json"
    if os.path.exists(middle_path):
        with open(middle_path, encoding="utf-8") as f:
            return json.load(f)
    else:
        with open(structure_path, encoding="utf-8") as f:
            return json.load(f)



