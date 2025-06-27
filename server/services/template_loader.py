from pptx.util import Pt, RGBColor
from pptx.enum.text import PP_ALIGN

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def set_shape_text(shape, text, font_size, font_color, font_name, align):
    shape.text = text
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.name = font_name
            run.font.color.rgb = RGBColor(*hex_to_rgb(font_color))
        if align == "center":
            paragraph.alignment = PP_ALIGN.CENTER
        elif align == "right":
            paragraph.alignment = PP_ALIGN.RIGHT
        else:
            paragraph.alignment = PP_ALIGN.LEFT

def fill_slide_text(slide, text, config):
    shape = slide.placeholders[config['placeholder_idx']]
    shape.text = text
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(config['font_size'])
            run.font.name = config['font_name']
            rgb = hex_to_rgb(config['font_color'])
            run.font.color.rgb = RGBColor(*rgb)

def fill_title_slide(slide, structure, config):
    title_text = next((item['text'] for item in structure if item.get('type') == 'title'), 'PPT标题')
    shape = slide.placeholders[config['placeholder_idx']]
    set_shape_text(shape, title_text, config['font_size'], config['font_color'], config['font_name'], config['align'])

def fill_toc_slide(slide, structure, config):
    toc_titles = [item['text'] for item in structure if item.get('type') == 'title' and item.get('level', 1) == 1]
    shape = slide.placeholders[config['placeholder_idx']]
    shape.text = ""
    for idx, item in enumerate(toc_titles, 1):
        p = shape.text_frame.add_paragraph()
        p.text = f"{idx}. {item}"
        p.level = 0
        p.font.size = Pt(config['font_size'])
        p.font.name = config['font_name']
        p.font.color.rgb = RGBColor(*hex_to_rgb(config['font_color']))
        p.alignment = getattr(PP_ALIGN, config['align'].upper(), PP_ALIGN.LEFT)

def fill_section_slide(slide, section_title, config):
    shape = slide.placeholders[config['placeholder_idx']]
    set_shape_text(shape, section_title, config['font_size'], config['font_color'], config['font_name'], config['align'])

def fill_content_slide(slide, title, content, config):
    title_shape = slide.placeholders[config['title_placeholder_idx']]
    content_shape = slide.placeholders[config['content_placeholder_idx']]
    set_shape_text(title_shape, title, config['title_font_size'], config['title_font_color'], config['font_name'], config['align'])
    set_shape_text(content_shape, content, config['content_font_size'], config['content_font_color'], config['font_name'], config['align'])

def fill_end_slide(slide, end_text, config):
    shape = slide.placeholders[config['placeholder_idx']]
    set_shape_text(shape, end_text, config['font_size'], config['font_color'], config['font_name'], config['align'])

# 目录、章节、内容、结束页同理
